import os
import pandas as pd
import numpy as np
import warnings
import sys
import re
from datetime import datetime
import gc

# --- VÉRIFICATION DES BIBLIOTHÈQUES ---
try:
    from pydeseq2.dds import DeseqDataSet
    from pydeseq2.ds import DeseqStats
except ImportError:
    print("\nERREUR : PyDESeq2 n'est pas détecté. Lancez : pip install pydeseq2")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("\nERREUR : python-pptx n'est pas détecté. Lancez : pip install python-pptx")
    sys.exit(1)

# --- CONFIGURATION DES CHEMINS ---
INPUT_DIR = os.environ.get("OMICS_IN_DIR", r"D:\Stage-CRBS\Stage_analyses\Transcriptomique")
METADATA_PATH = os.environ.get("OMICS_META_PATH", os.path.join(INPUT_DIR, "metadata.txt"))
OUTPUT_DIR = r"D:\Stage-CRBS\Stage_analyses\Résultats\Transcriptomique\Deseq2"

warnings.filterwarnings("ignore")

COHORT_MAPPING = {
    'A': "GSE160079_Raw_gene_counts_matrix_Cohort_MTM1-a.txt",
    'B': "GSE160081_Raw_gene_counts_matrix_Cohort_MTM1-b.txt",
    'C': "GSE160083_Raw_gene_counts_matrix_Cohort_MTM1-c.txt",
    'D': "GSE160077_Raw_gene_counts_matrix_Cohort_BIN1.txt",
    'E': "GSE160078_Raw_gene_counts_matrix_Cohort_DNM2_Updated-07-01-2024.xlsx",
    'F': "GSE282489_raw_counts_bin1_cohort.txt",
    'G': "GSE282489_raw_counts_dnm2_cohort.txt"
}

def normalize_name(name):
    return re.sub(r'[^a-zA-Z0-9]', '', str(name)).lower()

def load_metadata():
    if not os.path.exists(METADATA_PATH):
        raise FileNotFoundError(f"Metadata non trouvé: {METADATA_PATH}")
    df = pd.read_csv(METADATA_PATH, sep="\t")
    df['Sample_name'] = df['Sample_name'].astype(str).str.strip()
    return df

def calculate_cpm_means(counts_df, metadata_df):
    cpms = counts_df.div(counts_df.sum(axis=0), axis=1) * 1e6
    means = {}
    for group in ['WT', 'WT treated', 'Disease', 'Rescue']:
        samples = metadata_df[metadata_df['Group'] == group]['Sample_name'].tolist()
        valid_samples = [s for s in samples if s in cpms.columns]
        col_name = f'Moyenne_CPM_{group.replace(" ", "_")}'
        if valid_samples:
            means[col_name] = cpms[valid_samples].mean(axis=1)
        else:
            means[col_name] = 0.0
    return pd.DataFrame(means, index=counts_df.index)

def run_deseq2_logic(counts_df, metadata_df, contrast_num, contrast_den):
    sub_meta = metadata_df[metadata_df['Group'].isin([contrast_num, contrast_den])].copy()
    sub_meta.set_index('Sample_name', inplace=True)
    available = [c for c in sub_meta.index if c in counts_df.columns]
    if len(available) < 2: return None
    sub_meta = sub_meta.loc[available]
    sub_counts = counts_df[available].T
    sub_counts = sub_counts.loc[:, (sub_counts.sum(axis=0) > 0)]
    if sub_counts.empty or len(sub_meta['Group'].unique()) < 2: return None
    try:
        dds = DeseqDataSet(counts=sub_counts, metadata=sub_meta, design_factors="Group", quiet=True)
        dds.deseq2()
        stat_res = DeseqStats(dds, contrast=["Group", contrast_num, contrast_den], quiet=True)
        stat_res.summary() 
        return stat_res.results_df if hasattr(stat_res, 'results_df') else stat_res.results
    except Exception as e:
        print(f"   [Info] Échec DESeq2 ({contrast_num} vs {contrast_den}): {e}")
        return None
    finally:
        if 'dds' in locals(): del dds
        gc.collect()

def create_global_pptx(all_stats):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Rapport d'Analyse Transcriptomique (DESeq2)"
    slide.placeholders[1].text = f"Bilan Multi-Cohortes\nDate : {datetime.now().strftime('%d/%m/%Y')}"
    for stat in all_stats:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Synthèse Cohorte {stat['id']}"
        data = [
            ["Métrique / Indicateur", "Valeur / Résultat"],
            ["Total Gènes Analysés", f"{stat['n_genes']:,}"],
            ["Effectifs par Groupe (N)", f"WT: {stat['n_wt']} | WT_t: {stat['n_wtt']} | Dis: {stat['n_dis']} | Res: {stat['n_res']}"],
            ["Seuils Statistiques", "|Log2FC| > 1 et p-adj < 0.05"],
            ["Référence de Calibration", f"Groupe {stat['ref_grp']}"],
            ["Taille Signature Patho", f"{stat['n_patho_signif']} gènes"],
            ["Taille Efficacité Rescue", f"{stat['n_rescue_signif']} gènes corrigés"]
        ]
        rows, cols = len(data), 2
        table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(4)).table
        for r in range(rows):
            for c in range(cols): table.cell(r, c).text = data[r][c]
    prs.save(os.path.join(OUTPUT_DIR, "Deseq2_Rapport_Analyse_Transcriptomique.pptx"))

def process_cohort(cohort_id, metadata_all):
    file_name = COHORT_MAPPING.get(cohort_id)
    if not file_name: return None
    file_path = os.path.join(INPUT_DIR, file_name)
    if not os.path.exists(file_path): return None

    print(f"--- Analyse Cohorte {cohort_id} (Deseq2) ---")
    if file_name.endswith('.xlsx'): raw_counts = pd.read_excel(file_path, skiprows=1)
    else: raw_counts = pd.read_csv(file_path, sep=None, engine='python')

    raw_counts.columns = [str(c).strip() for c in raw_counts.columns]
    gene_col = raw_counts.columns[0]
    raw_counts.rename(columns={gene_col: 'Ensembl_Gene_ID'}, inplace=True)
    cohort_meta = metadata_all[metadata_all['Cohort'] == cohort_id].copy()
    mapping = {c: c for c in raw_counts.columns if c in cohort_meta['Sample_name'].values}
    if not mapping:
        norm_available = {normalize_name(c): c for c in raw_counts.columns if c != 'Ensembl_Gene_ID'}
        for sample in cohort_meta['Sample_name']:
            if normalize_name(sample) in norm_available: mapping[norm_available[normalize_name(sample)]] = sample

    raw_counts.rename(columns=mapping, inplace=True)
    raw_counts = raw_counts[['Ensembl_Gene_ID'] + list(mapping.values())]
    raw_counts.set_index('Ensembl_Gene_ID', inplace=True)
    raw_counts = raw_counts.apply(pd.to_numeric, errors='coerce').fillna(0).astype(int)

    cpm_df = calculate_cpm_means(raw_counts, cohort_meta)
    ref_grp = 'WT'
    seuil_1_pourcent = int(len(raw_counts) * 0.01)

    if 'WT treated' in cohort_meta['Group'].unique():
        res_calib = run_deseq2_logic(raw_counts, cohort_meta, 'WT treated', 'WT')
        if res_calib is not None:
            n_diff = (res_calib['padj'] < 0.05).sum()
            if n_diff < seuil_1_pourcent: ref_grp = 'WT treated'

    res_patho = run_deseq2_logic(raw_counts, cohort_meta, 'Disease', ref_grp)
    res_rescue = None
    if 'Rescue' in cohort_meta['Group'].unique(): res_rescue = run_deseq2_logic(raw_counts, cohort_meta, 'Rescue', 'Disease')

    final = pd.DataFrame(index=raw_counts.index)
    final['Pvalue_Patho'] = res_patho['padj'].reindex(final.index).fillna(1.0) if res_patho is not None else 1.0
    final['Log2FC_Patho'] = res_patho['log2FoldChange'].reindex(final.index).fillna(0.0) if res_patho is not None else 0.0
    final['Pvalue_Rescue'] = res_rescue['padj'].reindex(final.index).fillna(1.0) if res_rescue is not None else 1.0
    final['Log2FC_Rescue'] = res_rescue['log2FoldChange'].reindex(final.index).fillna(0.0) if res_rescue is not None else 0.0
    
    is_patho = (final['Pvalue_Patho'] < 0.05) & (final['Log2FC_Patho'].abs() > 1)
    opposed = np.sign(final['Log2FC_Patho']) != np.sign(final['Log2FC_Rescue'])
    is_rescue = (final['Pvalue_Rescue'] < 0.05) & is_patho & opposed

    # Noms de fichiers uniformisés
    final.to_csv(os.path.join(OUTPUT_DIR, f"Deseq2_full_analysis_Cohort_{cohort_id}.csv"), sep=";")
    final[is_patho].to_csv(os.path.join(OUTPUT_DIR, f"Deseq2_Genes_Signature_Patho_Cohort_{cohort_id}.csv"), sep=";")
    final[is_rescue].to_csv(os.path.join(OUTPUT_DIR, f"Deseq2_Genes_Efficacite_Rescue_Cohort_{cohort_id}.csv"), sep=";")
    
    grp_sizes = {
        'WT': len(cohort_meta[cohort_meta['Group'] == 'WT']),
        'WT treated': len(cohort_meta[cohort_meta['Group'] == 'WT treated']),
        'Disease': len(cohort_meta[cohort_meta['Group'] == 'Disease']),
        'Rescue': len(cohort_meta[cohort_meta['Group'] == 'Rescue'])
    }
    
    return {
        'id': cohort_id, 
        'ref_grp': ref_grp, 
        'n_genes': len(raw_counts),
        'n_wt': grp_sizes['WT'],
        'n_wtt': grp_sizes['WT treated'],
        'n_dis': grp_sizes['Disease'],
        'n_res': grp_sizes['Rescue'],
        'n_patho_signif': is_patho.sum(), 
        'n_rescue_signif': is_rescue.sum()
    }

def main():
    if not os.path.exists(OUTPUT_DIR): os.makedirs(OUTPUT_DIR)
    all_cohort_stats = []
    meta_data = load_metadata()
    for cid in ['A','B','C','D','E','F','G']:
        stats = process_cohort(cid, meta_data)
        if stats: all_cohort_stats.append(stats)
    if all_cohort_stats: create_global_pptx(all_cohort_stats)

if __name__ == "__main__":
    main()