import os
import pandas as pd
import numpy as np
import scipy.stats as stats
from statsmodels.stats.multitest import multipletests # Ajout pour p-adj
import warnings
from datetime import datetime

# --- CONFIGURATION DES CHEMINS ---
INPUT_DIR = r"D:\Stage-CRBS\Stage_analyses\Protéomique"
METADATA_PATH = os.path.join(INPUT_DIR, "metadata.tsv")
PROTEINS_PATH = os.path.join(INPUT_DIR, "proteinGroups.tsv")
OUTPUT_DIR = r"D:\Stage-CRBS\Stage_analyses\Résultats\Protéomique"

warnings.filterwarnings("ignore")

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

def clean_proteomics_data(df):
    cols_to_check = ['Potential contaminant', 'Reverse', 'Only identified by site']
    for col in cols_to_check:
        if col in df.columns:
            df = df[df[col] != '+']
    return df

def impute_missing_data(data):
    """
    Simule l'imputation par défaut de Perseus/wrProteo :
    Remplace les NaN par des valeurs issues d'une distribution normale décalée.
    """
    mu = data.mean()
    sigma = data.std()
    
    imputed_data = data.copy()
    for col in data.columns:
        n_nan = data[col].isnull().sum()
        if n_nan > 0:
            impute_mean = mu[col] - (1.8 * sigma[col])
            impute_std = sigma[col] * 0.3
            values = np.random.normal(loc=impute_mean, scale=impute_std, size=n_nan)
            imputed_data.loc[data[col].isnull(), col] = values
            
    return imputed_data

def calculate_stats_refined(df, group_test_cols, group_ref_cols):
    """
    Pipeline inspiré du papier : Log2 -> Imputation -> T-test -> BH Correction.
    """
    # 1. Passage en Log2 (remplace 0 par NaN)
    data_test = np.log2(df[group_test_cols].replace(0, np.nan))
    data_ref = np.log2(df[group_ref_cols].replace(0, np.nan))
    
    # 2. Imputation
    data_test_imp = impute_missing_data(data_test)
    data_ref_imp = impute_missing_data(data_ref)
    
    # 3. Log2FC (Test vs Ref)
    lfc = data_test_imp.mean(axis=1) - data_ref_imp.mean(axis=1)
    
    # 4. T-test indépendant
    pvals = df.apply(lambda row: stats.ttest_ind(
        data_test_imp.loc[row.name], 
        data_ref_imp.loc[row.name], 
        equal_var=True 
    ).pvalue, axis=1)
    
    # 5. Correction Benjamini-Hochberg (FDR)
    mask = ~pvals.isna()
    adj_pvals = pd.Series(1.0, index=pvals.index)
    if mask.any():
        _, corrected, _, _ = multipletests(pvals[mask], method='fdr_bh')
        adj_pvals[mask] = corrected
    
    return lfc, adj_pvals

def create_pptx_report(all_stats):
    if not PPTX_AVAILABLE: return
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Rapport d'Analyse Protéomique (MS/MS)"
    slide.placeholders[1].text = f"Analyse Différentielle MTM1-a (Pipeline MS/MS Count)\nDate : {datetime.now().strftime('%d/%m/%Y')}"

    for age, stat in all_stats.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Synthèse Age : {age} (MS/MS)"
        table_data = [
            ["Indicateur", "Valeur"],
            ["Protéines totales", str(stat['total'])],
            ["Signature Patho (|LFC|>1 & p-adj<0.05)", str(stat['n_patho'])],
            ["Protéines sauvées (Rescue)", str(stat['n_rescue'])],
            ["Top Gène sauvé", stat['top_rescue']]
        ]
        table = slide.shapes.add_table(5, 2, Inches(1), Inches(2), Inches(8), Inches(3)).table
        for r in range(5):
            for c in range(2):
                table.cell(r, c).text = table_data[r][c]
    prs.save(os.path.join(OUTPUT_DIR, "Rapport_Proteomique_MTM1a_ms.pptx"))

def main():
    if not os.path.exists(OUTPUT_DIR): os.makedirs(OUTPUT_DIR)
    
    print("Chargement des données pour analyse MS/MS...")
    df_meta = pd.read_csv(METADATA_PATH, sep='\t')
    df_prot = pd.read_csv(PROTEINS_PATH, sep='\t')
    
    df_prot = clean_proteomics_data(df_prot)
    df_prot.set_index('Protein IDs', inplace=True)
    
    final_df = df_prot[['Gene names', 'Protein names']].copy()
    report_stats = {}

    for age in df_meta['Age'].unique():
        print(f"--- Analyse Age : {age} (Mode MS/MS Count) ---")
        
        # Récupération des échantillons via le metadata
        samples_wt = df_meta[(df_meta['Age'] == age) & (df_meta['Genotype'] == 'WT')]['Sample'].tolist()
        samples_ko = df_meta[(df_meta['Age'] == age) & (df_meta['Genotype'] == 'KO')]['Sample'].tolist()
        samples_res = df_meta[(df_meta['Age'] == age) & (df_meta['Genotype'] == 'KO_Dnm2')]['Sample'].tolist()

        # Correspondance : on remplace 'LFQ intensity' par 'MS/MS count' pour pointer les bonnes colonnes
        cols_wt = [s.replace("LFQ intensity", "MS/MS count") for s in samples_wt]
        cols_ko = [s.replace("LFQ intensity", "MS/MS count") for s in samples_ko]
        cols_res = [s.replace("LFQ intensity", "MS/MS count") for s in samples_res]

        # Moyennes MS/MS pour affichage dans le CSV final
        final_df[f'Moyenne_MSMS_{age}_WT'] = df_prot[cols_wt].mean(axis=1)
        final_df[f'Moyenne_MSMS_{age}_KO'] = df_prot[cols_ko].mean(axis=1)
        final_df[f'Moyenne_MSMS_{age}_KO_Dnm2'] = df_prot[cols_res].mean(axis=1)

        # Stats Patho : KO (Test) vs WT (Ref)
        lfc_p, padj_p = calculate_stats_refined(df_prot, cols_ko, cols_wt)
        final_df[f'Log2FC_Patho_{age}'] = lfc_p
        final_df[f'Pvalue_Patho_{age}'] = padj_p 

        # Stats Rescue : KO_Dnm2 (Test) vs KO (Ref)
        lfc_r, padj_r = calculate_stats_refined(df_prot, cols_res, cols_ko)
        final_df[f'Log2FC_Rescue_{age}'] = lfc_r
        final_df[f'Pvalue_Rescue_{age}'] = padj_r

        # Seuils
        is_patho = (final_df[f'Pvalue_Patho_{age}'] < 0.05) & (final_df[f'Log2FC_Patho_{age}'].abs() > 1)
        opposed = np.sign(final_df[f'Log2FC_Patho_{age}']) != np.sign(final_df[f'Log2FC_Rescue_{age}'])
        is_rescue = (final_df[f'Pvalue_Rescue_{age}'] < 0.05) & is_patho & opposed

        final_df[f'Significatif_Patho_{age}'] = is_patho.map({True: 'OUI', False: 'NON'})
        final_df[f'Significatif_Rescue_{age}'] = is_rescue.map({True: 'OUI', False: 'NON'})

        # Exports avec suffixe _ms
        patho_df = final_df[final_df[f'Significatif_Patho_{age}'] == 'OUI']
        patho_df.to_csv(os.path.join(OUTPUT_DIR, f"Proteomics_Signature_Patho_{age}_ms.csv"), sep=";")
        
        rescue_df = final_df[final_df[f'Significatif_Rescue_{age}'] == 'OUI']
        rescue_df.to_csv(os.path.join(OUTPUT_DIR, f"Proteomics_Efficacite_Rescue_{age}_ms.csv"), sep=";")

        top_gene = rescue_df['Gene names'].iloc[0] if not rescue_df.empty else "Aucun"
        report_stats[age] = {
            'total': len(df_prot),
            'n_patho': is_patho.sum(),
            'n_rescue': is_rescue.sum(),
            'top_rescue': top_gene
        }

    final_df.to_csv(os.path.join(OUTPUT_DIR, "Proteomics_Analysis_Full_MTM1a_ms.csv"), sep=";")
    
    if PPTX_AVAILABLE:
        create_pptx_report(report_stats)
        print(f"Rapport PowerPoint généré : Rapport_Proteomique_MTM1a_ms.pptx")

    print(f"\nTerminé ! Résultats MS/MS générés dans : {OUTPUT_DIR}")

if __name__ == "__main__":
    main()