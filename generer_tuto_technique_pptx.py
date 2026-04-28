from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_tutorial_pptx():
    prs = Presentation()
    
    # helper: set dark background
    def set_dark_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(13, 17, 23) # Dark Navy

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_dark_bg(slide)
    title = slide.shapes.title
    title.text = "Guide de Déploiement OmicsFlow"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(88, 166, 255)
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Étape par étape : de votre PC local vers le Web\nAutomatisé avec GitHub & Streamlit"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(201, 209, 217)

    # Slide 2: Étape 1 - GitHub
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    title = slide.shapes.title
    title.text = "Étape 1 : Création du dépôt GitHub"
    content = slide.placeholders[1]
    content.text = ("1. Connectez-vous sur GitHub.com.\n"
                    "2. Cliquez sur [+] -> 'New repository'.\n"
                    "3. Nommez-le 'OmicsFlow-Studio'.\n"
                    "4. Choisir 'Public' (plus facile) ou 'Private'.\n"
                    "5. Cliquez sur 'Create repository'.")
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(201, 209, 217)
        p.font.size = Pt(22)

    # Slide 3: Étape 2 - Upload
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    title = slide.shapes.title
    title.text = "Étape 2 : Importation des fichiers"
    content = slide.placeholders[1]
    content.text = ("• Option Facile : Cliquez sur 'uploading an existing file' sur GitHub.\n"
                    "• Glissez-déposez TOUT le contenu du dossier local OmicsFlow_App.\n"
                    "• Option Pro : Initialisez un 'git push' depuis votre terminal :\n"
                    "  git init / git add . / git push origin main")
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(201, 209, 217)
        p.font.size = Pt(20)

    # Slide 4: Étape 3 - Streamlit Cloud
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    title = slide.shapes.title
    title.text = "Étape 3 : Streamlit Community Cloud"
    content = slide.placeholders[1]
    content.text = ("1. Allez sur share.streamlit.io.\n"
                    "2. Connectez votre compte GitHub.\n"
                    "3. Cliquez sur 'Create app'.\n"
                    "4. Sélectionnez votre dépôt et le fichier 'app.py'.\n"
                    "5. Cliquez sur 'Deploy!'")
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(201, 209, 217)
        p.font.size = Pt(22)

    # Slide 5: Behind the Screens
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    title = slide.shapes.title
    title.text = "Phase de Construction Automatique"
    content = slide.placeholders[1]
    content.text = ("• packages.txt : Streamlit installe R et les librairies système Debian.\n"
                    "• requirements.txt : Installation automatique de PyDeseq2, Plotly, etc.\n"
                    "• Temps estimé : 4 à 5 minutes pour le premier déploiement.\n"
                    "• URL Finale : https://omicsflow-studio-xxxx.streamlit.app")
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(201, 209, 217)
        p.font.size = Pt(18)

    save_path = r"D:\Stage-CRBS\Stage_analyses\OmicsFlow_App\Tutorial_Technique_Deploiement.pptx"
    prs.save(save_path)
    print(f"Tutorial PPTX saved to {save_path}")

if __name__ == "__main__":
    create_tutorial_pptx()
