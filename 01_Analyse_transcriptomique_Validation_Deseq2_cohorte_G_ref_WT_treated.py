import os
import pandas as pd
import numpy as np
import warnings
import sys
import re

# --- VÉRIFICATION DES BIBLIOTHÈQUES ---
try:
    from pydeseq2.dds import DeseqDataSet
    from pydeseq2.ds import DeseqStats
except ImportError:
    print("\nERREUR : PyDESeq2 non détecté. Lancez : pip install pydeseq2")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("\nERREUR : python-pptx non détecté. Lancez : pip install python-pptx")
    sys.exit(1)

# --- CONFIGURATION DES CHEMINS ---
METADATA_PATH = r"D:\Stage-CRBS\Stage_analyses\Transcriptomique\metadata.txt"
INPUT_DIR = r"D:\Stage-CRBS\Stage_analyses\Transcriptomique"
OUTPUT_DIR = r"D:\Stage-CRBS\Stage_analyses\Résultats\Transcriptomique"

# Paramètres spécifiques Cohorte G
COHORT_ID = 'G'
FILENAME_G = "GSE282489_raw_counts_dnm2_cohort.txt"
REFERENCE_FIXE = 'WT treated'

warnings.filterwarnings("ignore")

def normalize_name(name):
    """Normalise les noms pour la correspondance metadata/counts."""
    return re.sub(r'[^a-zA-Z0-9]', '', str(name)).lower()

def load_metadata():
    if not os.path.exists(METADATA_PATH):
        raise FileNotFoundError(f"Metadata non trouvé: {METADATA_PATH}")
    df = pd.read_csv(METADATA_PATH, sep="\t")
    df['Sample_name'] = df['Sample_name'].str.strip()
    return df

def run_deseq2_robust(counts_df, metadata_df, contrast_num, contrast_den):
    """Exécute DESeq2 de manière robuste."""
    sub_meta = metadata_df[metadata_df['Group'].isin([contrast_num, contrast_den])].copy()
    common_samples = [s for s in sub_meta['Sample_name'] if s in counts_df.columns]
    
    if len(common_samples) < 2:
        return None

    sub_meta = sub_meta[sub_meta['Sample_name'].isin(common_samples)].set_index('Sample_name')
    sub_counts = counts_df[common_samples].T
    
    # Filtrer les gènes avec uniquement des zéros
    sub_counts = sub_counts.loc[:, (sub_counts.sum(axis=0) > 0)]
    
    if len(sub_meta['Group'].unique()) < 2:
        return None

    try:
        dds = DeseqDataSet(counts=sub_counts, metadata=sub_meta, design_factors="Group", quiet=True)
        dds.deseq2()
        stat_res = DeseqStats(dds, contrast=["Group", contrast_num, contrast_den], quiet=True)
        stat_res.summary()
        return stat_res.results_df
    except Exception as e:
        print(f"   [Erreur DESeq2 {contrast_num} vs {contrast_den}] : {e}")
        return None

def create_pptx(stat):
    """Génère le rapport PowerPoint de synthèse."""
    prs = Presentation()
    
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Analyse Différentielle Cohorte {COHORT_ID}"
    slide.placeholders[1].text = f"Référence : {REFERENCE_FIXE}\nFichier : {FILENAME_G}"
    
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Synthèse des Gènes - Cohorte {COHORT_ID}"
    
    data = [
        ["Indicateur", "Valeur"],
        ["Signature Patho (vs WT treated)", str(stat['n_patho'])],
        ["Efficacité Rescue (vs Disease)", str(stat['n_rescue'])],
        ["Seuils", "P-adj < 0.05 | |Log2FC| > 1"]
    ]
    
    table = slide.shapes.add_table(4, 2, Inches(1), Inches(2), Inches(8), Inches(2)).table
    for r in range(4):
        for c in range(2):
            table.cell(r, c).text = data[r][c]

    filename = f"Deseq2_Rapport_G_ref_WT_treated.pptx"
    path = os.path.join(OUTPUT_DIR, filename)
    prs.save(path)
    return filename

def main():
    if not os.path.exists(OUTPUT_DIR): os.makedirs(OUTPUT_DIR)
    
    # 1. Chargement des métadonnées
    meta = load_metadata()
    cohort_meta = meta[meta['Cohort'] == COHORT_ID].copy()
    
    # 2. Chargement des counts
    file_path = os.path.join(INPUT_DIR, FILENAME_G)
    if not os.path.exists(file_path):
        print(f"ERREUR : Fichier introuvable au chemin : {file_path}")
        return

    print(f"Chargement des counts : {FILENAME_G}")
    raw_counts = pd.read_csv(file_path, sep=None, engine='python')
    raw_counts.columns = [str(c).strip() for c in raw_counts.columns]
    gene_col = raw_counts.columns[0]
    raw_counts.set_index(gene_col, inplace=True)

    # 3. Alignement des noms d'échantillons
    mapping = {}
    cols_norm = {normalize_name(c): c for c in raw_counts.columns}
    for sample in cohort_meta['Sample_name']:
        norm_s = normalize_name(sample)
        if norm_s in cols_norm:
            mapping[cols_norm[norm_s]] = sample
    
    raw_counts.rename(columns=mapping, inplace=True)
    valid_cols = [c for c in raw_counts.columns if c in cohort_meta['Sample_name'].values]
    
    if not valid_cols:
        print("Erreur : Aucun échantillon de la cohorte G n'a été trouvé dans le fichier de counts.")
        return

    counts = raw_counts[valid_cols].apply(pd.to_numeric, errors='coerce').fillna(0).astype(int)
    print(f"Échantillons identifiés pour analyse : {len(valid_cols)}")

    # 4. Analyses DESeq2
    print(f"Calcul Pathologie : Disease vs {REFERENCE_FIXE}...")
    res_patho = run_deseq2_robust(counts, cohort_meta, 'Disease', REFERENCE_FIXE)
    
    print("Calcul Sauvetage : Rescue vs Disease...")
    res_rescue = run_deseq2_robust(counts, cohort_meta, 'Rescue', 'Disease')

    # 5. Compilation et filtrage
    final = pd.DataFrame(index=counts.index)
    if res_patho is not None:
        final['Log2FC_Patho'] = res_patho['log2FoldChange']
        final['Pvalue_Patho'] = res_patho['padj']
    else:
        final['Log2FC_Patho'], final['Pvalue_Patho'] = 0, 1
        
    if res_rescue is not None:
        final['Log2FC_Rescue'] = res_rescue['log2FoldChange']
        final['Pvalue_Rescue'] = res_rescue['padj']
    else:
        final['Log2FC_Rescue'], final['Pvalue_Rescue'] = 0, 1

    is_patho = (final['Pvalue_Patho'] < 0.05) & (final['Log2FC_Patho'].abs() > 1)
    is_rescue = (final['Pvalue_Rescue'] < 0.05) & is_patho & (np.sign(final['Log2FC_Patho']) != np.sign(final['Log2FC_Rescue']))

    # --- CHANGEMENT DEMANDÉ : NOMMER L'INDEX POUR L'EXPORT ---
    final.index.name = "Ensembl_Gene_ID"

    # 6. Exports
    tag = f"Cohort_G_ref_WT_treated"
    final.to_csv(os.path.join(OUTPUT_DIR, f"Deseq2_Full_Results_{tag}.csv"), sep=";")
    
    patho_df = final[is_patho]
    patho_df.to_csv(os.path.join(OUTPUT_DIR, f"Deseq2_Genes_Signature_Patho_{tag}.csv"), sep=";")
    
    rescue_df = final[is_rescue]
    rescue_df.to_csv(os.path.join(OUTPUT_DIR, f"Deseq2_Genes_Efficacite_Rescue_{tag}.csv"), sep=";")

    print(f"\nRésultats pour la Cohorte G :")
    print(f" - Signature Patho (vs {REFERENCE_FIXE}) : {len(patho_df)} gènes")
    print(f" - Gènes Rescue (vs Disease) : {len(rescue_df)} gènes")
    
    pptx_file = create_pptx({'n_patho': len(patho_df), 'n_rescue': len(rescue_df)})
    print(f"Rapport enregistré : {pptx_file}")

if __name__ == "__main__":
    main()