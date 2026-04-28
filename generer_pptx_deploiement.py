from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_deployment_pptx():
    prs = Presentation()
    
    # helper: set dark background
    def set_dark_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(13, 17, 23) # Dark Navy

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_dark_bg(slide)
    title = slide.shapes.title
    title.text = "OmicsFlow Studio : Déploiement Cloud"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(88, 166, 255)
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Migration vers une plateforme interactive & portable\nAvril 2026"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(201, 209, 217)

    # 2. Architecture Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    title = slide.shapes.title
    title.text = "Architecture Isolation & Sécurité"
    content = slide.placeholders[1]
    content.text = ("• Sessions Isolées : Création de dossiers session uniques pour chaque run.\n"
                    "• Protection des Données : Aucun écrasement des fichiers sources.\n"
                    "• Logic Locks : Pipelines bloqués jusqu'à l'upload effectif des données.")
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(201, 209, 217)
        p.font.size = Pt(20)

    # 3. Portabilité Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    title = slide.shapes.title
    title.text = "Cloud-Ready : Zéro Dépendance Locale"
    content = slide.placeholders[1]
    content.text = ("• Refactorisation : Remplacement des chemins D:\\ par des chemins relatifs.\n"
                    "• OmicsFlow_App : Un dossier 'clône' dédié uniquement au Web.\n"
                    "• Config : requirements.txt (Python) + packages.txt (R System).")
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(201, 209, 217)
        p.font.size = Pt(20)

    # 4. Multi-Omics Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    title = slide.shapes.title
    title.text = "Pipeline Unifié"
    content = slide.placeholders[1]
    content.text = ("• Transcriptomique : PyDeseq2 + PCA dynamique.\n"
                    "• Protéomique : Double moteur Welch (Pingouin) et DEP (R/Limma).\n"
                    "• Métabolomique : Welch + Correction FDR Benjamini-Hochberg.\n"
                    "• Visualisation : Volcano Plots interactifs Plotly.")
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(201, 209, 217)
        p.font.size = Pt(20)

    # 5. Continuous Deployment Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    title = slide.shapes.title
    title.text = "Workflow de Mise à Jour"
    content = slide.placeholders[1]
    content.text = ("1. Modification dans le dossier local OmicsFlow_App.\n"
                    "2. Envoi vers GitHub (Push).\n"
                    "3. Mise à jour automatique sur Streamlit Cloud.\n"
                    "🚀 L'appli Web vit séparément de vos fichiers de recherche.")
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(201, 209, 217)
        p.font.size = Pt(20)

    save_path = r"D:\Stage-CRBS\Stage_analyses\Résultats\Deployment_Report_OmicsFlow.pptx"
    prs.save(save_path)
    print(f"PPTX saved to {save_path}")

if __name__ == "__main__":
    create_deployment_pptx()
