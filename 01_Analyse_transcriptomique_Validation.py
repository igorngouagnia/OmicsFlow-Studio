import os
import pandas as pd
import numpy as np
from scipy.stats import ttest_ind
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# --- CONFIGURATION DES CHEMINS ---
METADATA_PATH = r"D:\Stage-CRBS\Stage_analyses\Transcriptomique\metadata.txt"
INPUT_DIR = r"D:\Stage-CRBS\Stage_analyses\Transcriptomique"
OUTPUT_DIR = r"D:\Stage-CRBS\Stage_analyses\Résultats\Transcriptomique"
LOG_PATH = os.path.join(OUTPUT_DIR, "log_ARN_Synthese.txt")
PPTX_PATH = os.path.join(OUTPUT_DIR, "Rapport_Analyse_Transcriptomique.pptx")

# Mapping des identifiants de cohorte vers les fichiers physiques fournis par Julie
COHORT_MAPPING = {
    'A': "GSE160079_Raw_gene_counts_matrix_Cohort_MTM1-a.txt",
    'B': "GSE160081_Raw_gene_counts_matrix_Cohort_MTM1-b.txt",
    'C': "GSE160083_Raw_gene_counts_matrix_Cohort_MTM1-c.txt",
    'D': "GSE160077_Raw_gene_counts_matrix_Cohort_BIN1.txt",
    'E': "GSE160078_Raw_gene_counts_matrix_Cohort_DNM2_Updated-07-01-2024.xlsx",
    'F': "GSE282489_raw_counts_bin1_cohort.txt",
    'G': "GSE282489_raw_counts_dnm2_cohort.txt"
}

def load_metadata():
    """Charge le fichier de correspondance échantillons/groupes."""
    if not os.path.exists(METADATA_PATH):
        raise FileNotFoundError(f"Metadata non trouvé: {METADATA_PATH}")
    df = pd.read_csv(METADATA_PATH, sep="\t")
    df['Sample_name'] = df['Sample_name'].str.strip()
    return df

def calculate_cpm(counts_df):
    """
    Normalisation CPM (Counts Per Million).
    Indispensable pour comparer des échantillons qui n'ont pas la même profondeur de séquençage.
    """
    gene_ids = counts_df.iloc[:, 0]
    numeric_data = counts_df.iloc[:, 1:]
    # Somme des lectures par échantillon
    totals = numeric_data.sum(axis=0).replace(0, 1)
    # Division par le total et multiplication par 1 million
    cpm_df = (numeric_data / totals) * 1_000_000
    return pd.concat([gene_ids, cpm_df], axis=1)

def perform_ttest(df, grp1_samples, grp2_samples):
    """
    Test t de Welch (equal_var=False).
    Adapté à la biologie car il n'assume pas que les groupes (ex: Sain vs Malade) 
    ont la même variabilité (variance).
    """
    if len(grp1_samples) < 2 or len(grp2_samples) < 2:
        return np.ones(len(df)) # Retourne p-value de 1.0 si effectifs insuffisants
    
    data1 = df[grp1_samples].values
    data2 = df[grp2_samples].values
    
    # Calcul statistique sur l'axe des gènes
    _, p_val = ttest_ind(data1, data2, axis=1, equal_var=False, nan_policy='omit')
    return np.nan_to_num(p_val, nan=1.0)

def create_pptx_template():
    """Initialise le support de présentation PowerPoint."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Rapport d'Analyse Transcriptomique"
    slide.placeholders[1].text = "Analyse Multi-Cohortes : Pathologie & Effet Rescue\nDate : 03/03/2026"
    return prs

def add_cohort_slide(prs, cohort_id, stats):
    """Génère une diapositive de résumé pour chaque cohorte."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Synthèse Cohorte {cohort_id}"
    
    # On passe à 8 lignes pour intégrer le ratio_patho
    rows, cols = 8, 2
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)).table
    
    data = [
        ("Indicateur", "Valeur"),
        ("Échantillons WT / WT Treated", f"{stats['n_WT']} / {stats['n_WT_treated']}"),
        ("Échantillons Disease / Rescue", f"{stats['n_Disease']} / {stats['n_Rescue']}"),
        ("Référence Choisie", stats['ref']),
        ("Signature Patho (|L2FC|>1 & P<0.05)", str(stats['sig_pathog'])),
        ("Ratio Patho / Total Gènes", f"{stats['ratio_patho']:.2%}"),
        ("Efficacité Rescue (Patho Signif. + Rescue Signif.)", str(stats['sig_rescue'])),
        ("Ratio Rescue / Patho", f"{stats['ratio_sig']:.1%}")
    ]
    
    for i, (label, val) in enumerate(data):
        table.cell(i, 0).text = label
        table.cell(i, 1).text = val
        for cell in [table.cell(i, 0), table.cell(i, 1)]:
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(15)
            para.alignment = PP_ALIGN.CENTER if cell == table.cell(i, 1) else PP_ALIGN.LEFT

def process_cohort(cohort_id, metadata_df, log_file, prs):
    """Fonction principale de traitement d'une cohorte."""
    file_name = COHORT_MAPPING.get(cohort_id)
    if not file_name: return
    file_path = os.path.join(INPUT_DIR, file_name)
    
    if not os.path.exists(file_path):
        print(f"Fichier manquant : {file_name}")
        return

    # 1. CHARGEMENT (Gestion Multi-format)
    if file_name.endswith('.xlsx'):
        counts = pd.read_excel(file_path, skiprows=1)
    else:
        counts = pd.read_csv(file_path, sep="\t")
    
    # Nettoyage des noms de colonnes
    counts.columns = [str(c).strip() for c in counts.columns]
    id_col = counts.columns[0]
    counts = counts.rename(columns={id_col: 'Ensembl_Gene_ID'})

    # 2. MATCHING ÉCHANTILLONS (Logique de nettoyage textuel)
    # On filtre les métadonnées pour la cohorte actuelle
    cohort_meta = metadata_df[metadata_df['Cohort'] == cohort_id]
    valid_samples = cohort_meta['Sample_name'].tolist()
    
    # Fonction pour ignorer espaces, tirets et majuscules lors du matching
    def clean(s): return "".join(filter(str.isalnum, str(s))).lower()
    norm_map = {clean(s): s for s in valid_samples}
    
    rename_cols = {}
    keep_cols = ['Ensembl_Gene_ID']
    for c in counts.columns:
        if c == 'Ensembl_Gene_ID': continue
        c_clean = clean(c)
        if c_clean in norm_map:
            m_n = norm_map[c_clean]
            rename_cols[c] = m_n
            keep_cols.append(m_n)
            
    counts = counts.rename(columns=rename_cols)[keep_cols]
    if counts.shape[1] < 2: return

    # 3. NORMALISATION & CALCUL DES MOYENNES
    cpm = calculate_cpm(counts)
    group_map = {
        'WT': cohort_meta[cohort_meta['Group'] == 'WT']['Sample_name'].tolist(),
        'WT_treated': cohort_meta[cohort_meta['Group'] == 'WT treated']['Sample_name'].tolist(),
        'Disease': cohort_meta[cohort_meta['Group'] == 'Disease']['Sample_name'].tolist(),
        'Rescue': cohort_meta[cohort_meta['Group'] == 'Rescue']['Sample_name'].tolist()
    }
    # On ne garde que les échantillons réellement présents dans le fichier de données
    grp_c = {k: [s for s in v if s in cpm.columns] for k, v in group_map.items()}
    
    res = pd.DataFrame({'Ensembl_Gene_ID': cpm['Ensembl_Gene_ID']})
    for k, v in grp_c.items():
        res[f'Moyenne_CPM_{k}'] = cpm[v].mean(axis=1) if v else 0.0

    # 4. CALIBRATION DU CONTRÔLE (Choix WT vs WT_treated)
    # On vérifie si le traitement seul impacte le groupe Sain (Sain Traité vs Sain Pur)
    ratio = 1.0
    if grp_c['WT'] and grp_c['WT_treated']:
        # +0.1 est le pseudocount pour éviter les erreurs mathématiques (division par zéro)
        cal_l2fc = np.log2((res['Moyenne_CPM_WT_treated'] + 0.1) / (res['Moyenne_CPM_WT'] + 0.1))
        # Si moins de 1% des gènes bougent significativement, on considère le traitement "neutre"
        ratio = (cal_l2fc.abs() > 1).sum() / len(res) if len(res) > 0 else 1.0

    # Verdict : Si neutre (<1%), la référence devient le WT_treated pour être iso-traitement.
    if ratio < 0.01 and grp_c['WT_treated']:
        ref_grp, ref_col, ref_samples = 'WT_treated', 'Moyenne_CPM_WT_treated', grp_c['WT_treated']
    else:
        ref_grp, ref_col, ref_samples = 'WT', 'Moyenne_CPM_WT', grp_c['WT']

    # 5. ANALYSE STATISTIQUE DES DIFFÉRENCES
    # Signature de la Maladie : Malade vs Référence saine choisie
    res['Pvalue_Patho'] = perform_ttest(cpm, grp_c['Disease'], ref_samples)
    # Effet du traitement (Rescue) : Rescue vs Malade
    res['Pvalue_Rescue'] = perform_ttest(cpm, grp_c['Rescue'], grp_c['Disease'])
    
    res['Ref_Utilisee'] = ref_grp
    res['Log2FC_Patho'] = np.log2((res['Moyenne_CPM_Disease'] + 0.1) / (res[ref_col] + 0.1))
    res['Log2FC_Rescue'] = np.log2((res['Moyenne_CPM_Rescue'] + 0.1) / (res['Moyenne_CPM_Disease'] + 0.1))

    # Marquage des gènes significatifs (Seuil Log2FC > 1 et p < 0.05)
    res['Significatif_Patho'] = ((res['Log2FC_Patho'].abs() > 1) & (res['Pvalue_Patho'] < 0.05)).map({True:'OUI', False:'NON'})
    res['Significatif_Rescue'] = ((res['Log2FC_Rescue'].abs() > 1) & (res['Pvalue_Rescue'] < 0.05)).map({True:'OUI', False:'NON'})

    # 6. EXPORT DES 3 FICHIERS CIBLES
    # Fichier Global : Tous les gènes analysés
    res.to_csv(os.path.join(OUTPUT_DIR, f"Genes_Analysis_Cohort_{cohort_id}.csv"), index=False, sep=";")
    
    # Fichier Signature Patho : Uniquement les gènes définissant la maladie
    patho_df = res[res['Significatif_Patho'] == 'OUI']
    patho_df.to_csv(os.path.join(OUTPUT_DIR, f"Genes_Signature_Patho_Cohort_{cohort_id}.csv"), index=False, sep=";")
    
    # Fichier Efficacité Rescue : Gènes Patho corrigés par le traitement (double significativité)
    rescue_df = res[(res['Significatif_Patho'] == 'OUI') & (res['Significatif_Rescue'] == 'OUI')].copy()
    # Calcul de la puissance de correction (en %)
    rescue_df['Pourcentage_Correction'] = (rescue_df['Log2FC_Rescue'].abs() / (rescue_df['Log2FC_Patho'].abs() + 1e-9)) * 100
    rescue_df.to_csv(os.path.join(OUTPUT_DIR, f"Genes_Efficacite_Rescue_Cohort_{cohort_id}.csv"), index=False, sep=";")

    # 7. LOG ET PRÉSENTATION
    n_p, n_r = len(patho_df), len(rescue_df)
    log_file.write(f"\nCOHORTE {cohort_id} : Patho={n_p}, Rescue={n_r}, Ref={ref_grp}\n")
    
    add_cohort_slide(prs, cohort_id, {
        'n_WT': len(grp_c['WT']), 'n_WT_treated': len(grp_c['WT_treated']),
        'n_Disease': len(grp_c['Disease']), 'n_Rescue': len(grp_c['Rescue']),
        'ref': ref_grp, 
        'sig_pathog': n_p, 
        'ratio_patho': n_p / len(res) if len(res) > 0 else 0,
        'sig_rescue': n_r,
        'ratio_sig': n_r / n_p if n_p > 0 else 0
    })

def main():
    """Point d'entrée du script."""
    if not os.path.exists(OUTPUT_DIR): os.makedirs(OUTPUT_DIR)
    meta = load_metadata()
    prs = create_pptx_template()
    
    with open(LOG_PATH, "w", encoding="utf-8") as log_file:
        log_file.write("=== LOG ANALYSE TRANSCRIPTOMIQUE V3 ===\n")
        # Boucle sur les 7 cohortes de l'étude
        for cid in ['A', 'B', 'C', 'D', 'E', 'F', 'G']:
            print(f"Traitement Cohorte {cid}...")
            try:
                process_cohort(cid, meta, log_file, prs)
            except Exception as e:
                print(f"Erreur critique sur cohorte {cid}: {e}")
                log_file.write(f"ERREUR {cid}: {e}\n")
    
    prs.save(PPTX_PATH)
    print(f"Terminé. Les fichiers sont disponibles dans : {OUTPUT_DIR}")

if __name__ == "__main__":
    main()